"""Импорт реального датасета хакатона -> data/corpus.json.

Рекурсивно обходит папку, парсит PDF/DOCX/DOC/PPTX/XLS/XLSX/DOCM, распаковывает ZIP.
Определяет заголовок, происхождение (по языку), год, категорию (по верхней папке).
Пустые/нечитаемые (сканы без текстового слоя) отсеиваются.

Использование:
  python scripts/ingest_real.py "/path/to/Задача 2. Научный клубок" [--limit N]
"""
from __future__ import annotations

import json
import os
import re
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
MAX_CHARS = 9000        # сколько текста храним на документ
MAX_PDF_PAGES = 20      # ограничение по страницам PDF для скорости
MIN_CHARS = 200         # меньше — считаем нечитаемым (скан/пусто)


def _clean(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def parse_pdf(path: str) -> str:
    import fitz
    doc = fitz.open(path)
    parts = []
    for i, page in enumerate(doc):
        if i >= MAX_PDF_PAGES:
            break
        parts.append(page.get_text())
    doc.close()
    return _clean("\n".join(parts))


def parse_docx(path: str) -> str:
    from docx import Document
    d = Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for tbl in d.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return _clean("\n".join(parts))


def parse_doc(path: str) -> str:
    # старый .doc — через macOS textutil
    out = subprocess.run(["textutil", "-convert", "txt", "-stdout", path],
                         capture_output=True, timeout=60)
    return _clean(out.stdout.decode("utf-8", "ignore"))


def parse_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return _clean("\n".join(parts))


def parse_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets[:3]:
        for r, row in enumerate(ws.iter_rows(values_only=True)):
            if r > 200:
                break
            vals = [str(v) for v in row if v is not None]
            if vals:
                parts.append(" | ".join(vals))
    wb.close()
    return _clean("\n".join(parts))


def parse_xls(path: str) -> str:
    import xlrd
    book = xlrd.open_workbook(path)
    parts = []
    for sh in book.sheets()[:3]:
        for r in range(min(sh.nrows, 200)):
            vals = [str(sh.cell_value(r, c)) for c in range(sh.ncols)
                    if str(sh.cell_value(r, c)).strip()]
            if vals:
                parts.append(" | ".join(vals))
    return _clean("\n".join(parts))


PARSERS = {
    ".pdf": parse_pdf, ".docx": parse_docx, ".docm": parse_docx, ".doc": parse_doc,
    ".pptx": parse_pptx, ".xlsx": parse_xlsx, ".xls": parse_xls,
}


def detect_origin(text: str) -> str:
    cyr = len(re.findall(r"[а-яА-Я]", text))
    lat = len(re.findall(r"[a-zA-Z]", text))
    return "Зарубежная" if lat > cyr * 1.2 else "Отечественная"


def detect_year(text: str):
    ys = [int(y) for y in re.findall(r"\b((?:19|20)\d{2})\b", text[:MAX_CHARS])]
    ys = [y for y in ys if 1990 <= y <= 2025]
    return max(ys) if ys else None


def clean_filename(fn: str) -> str:
    base = os.path.splitext(os.path.basename(fn))[0]
    base = re.sub(r"^\d+[\s._-]*", "", base)  # ведущий номер
    return base.strip(" _.-")


def make_title(text: str, fn: str) -> str:
    for line in text.split("\n")[:8]:
        s = line.strip()
        if 15 <= len(s) <= 180 and s.count("–") < 2 and s.count("-") < 4 and not s.isupper():
            return s
    return clean_filename(fn) or os.path.basename(fn)


def category_of(path: str, base: str) -> str:
    rel = os.path.relpath(path, base)
    parts = rel.split(os.sep)
    # "Источники информации/<Категория>/..."
    if len(parts) >= 2 and parts[0].startswith("Источники"):
        return parts[1]
    return parts[0] if parts else ""


def iter_files(base: str, workdir: str):
    """Отдаёт (путь, категория). Распаковывает zip во временную папку."""
    for dirpath, _, files in os.walk(base):
        for f in files:
            p = os.path.join(dirpath, f)
            ext = os.path.splitext(f)[1].lower()
            if ext == ".zip":
                try:
                    dest = os.path.join(workdir, f + "_x")
                    with zipfile.ZipFile(p) as z:
                        z.extractall(dest)
                    cat = category_of(p, base)
                    for dp, _, fs in os.walk(dest):
                        for ff in fs:
                            if os.path.splitext(ff)[1].lower() in PARSERS:
                                yield os.path.join(dp, ff), cat
                except Exception:  # noqa: BLE001
                    continue
            elif ext in PARSERS:
                yield p, category_of(p, base)


def main():
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    base = args[0] if args else \
        "/Users/maximizyumov/Downloads/Задача 2. Научный клубок"
    limit = None
    if "--limit" in sys.argv:
        limit = int(sys.argv[sys.argv.index("--limit") + 1])

    docs, n_ok, n_skip, n_err = [], 0, 0, 0
    with tempfile.TemporaryDirectory() as workdir:
        files = list(iter_files(base, workdir))
        print(f"Найдено файлов для парсинга: {len(files)}")
        for i, (path, cat) in enumerate(sorted(files, key=lambda x: x[0])):
            ext = os.path.splitext(path)[1].lower()
            try:
                text = PARSERS[ext](path)
            except Exception as e:  # noqa: BLE001
                n_err += 1
                if n_err <= 10:
                    print(f"  [err] {os.path.basename(path)[:50]}: {str(e)[:60]}")
                continue
            if len(text) < MIN_CHARS:
                n_skip += 1
                continue
            n_ok += 1
            docs.append({
                "id": f"D-{n_ok:04d}",
                "title": make_title(text, path),
                "origin": detect_origin(text),
                "year": detect_year(text),
                "category": cat,
                "path": os.path.relpath(path, base),
                "text": text[:MAX_CHARS],
            })
            if n_ok % 100 == 0:
                print(f"  распарсено {n_ok}…")
            if limit and n_ok >= limit:
                break

    out = ROOT / "data" / "corpus.json"
    out.write_text(json.dumps(docs, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"\nГотово: {n_ok} документов (пропущено пустых/сканов: {n_skip}, ошибок: {n_err})")
    from collections import Counter
    print("По категориям:", dict(Counter(d["category"] for d in docs)))
    print("По происхождению:", dict(Counter(d["origin"] for d in docs)))
    print(f"-> {out}")


if __name__ == "__main__":
    main()
